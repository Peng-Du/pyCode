from google.cloud import translate_v3 as translate
import os
from pptx import Presentation

def translate_pptx(pptx_file):
    # 设置Google Cloud API密钥路径
    # os.environ['GOOGLE_APPLICATION_CREDENTIALS'] = ''
    # set GOOGLE_APPLICATION_CREDENTIALS=
	# echo "GOOGLE_APPLICATION_CREDENTIALS"

    # 创建Google Translate客户端
    client = translate.TranslationServiceClient()

    # 加载PPTX文件
    presentation = Presentation(pptx_file)

    # 创建翻译请求
    parent = "projects/transppt-123456/locations/global"
    target_language = "es"
    model = "projects/transppt-123456/locations/global/models/general/nmt"
    contents = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        contents.append(run.text)
                        # print(run.text)
    if not contents:
        print("No text content found in the PPTX file. Exiting translation.")
        return

    request = translate.TranslateTextRequest(
        parent=parent,
        target_language_code=target_language,
        contents=contents,
        model=model
    )

    response = client.translate_text(request=request)
    

    # 更新PPTX中的文本
    translated_text_index = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = response.translations[translated_text_index].translated_text
                        # print(run.text)
                        translated_text_index += 1

    # 保存处理后的PPT文件
    translated_pptx_file = pptx_file.split(".pptx")[0] + "_translated.pptx"
    presentation.save(translated_pptx_file)

    print("PPTX translation complete. Translated file saved as", translated_pptx_file)

if __name__ == '__main__':
    pptx_file = "test.pptx"
    translate_pptx(pptx_file)

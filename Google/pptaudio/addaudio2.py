import collections
import collections.abc
from pptx import Presentation
import docx
import os
import re
import aspose.slides as slides
import aspose.slides.animation as animation
import wave
import math
# need library: python-pptx, python-docx, azure-cognitiveservices-speech, aspose.slides, wave


input_pptx = "test.pptx"
split_parts  = input_pptx.rsplit('.', 1)
# 移除最后一部分
pptx_file_name = split_parts[0]

docx_file = pptx_file_name + ".docx"

def Add_Audio_to_PPTX():
    # load presentation
    with slides.Presentation(input_pptx) as presentation:
        i = 0
        for sld in presentation.slides:
            i += 1
            # load the wav sound file to stream
            audio_file = 'audio\\' + f'P{i:02d}.wav'
            if not os.path.isfile(audio_file):
                continue
            with open(audio_file, "rb") as in_file:
                # add audio frame
                audio_frame = sld.shapes.add_audio_frame_embedded(150, 100, 30, 30, in_file)

                # set play mode and volume of the audio
                audio_frame.play_mode = slides.AudioPlayModePreset.AUTO
                audio_frame.volume = slides.AudioVolumeMode.LOUD
                audio_frame.hide_at_showing = True
            #get time duration
            f = wave.open(audio_file, 'rb')
            time_count = f.getparams().nframes/f.getparams().framerate
            time_in_seconds = math.ceil(time_count) + 1
            #print(time_in_seconds)
            sld.slide_show_transition.advance_on_click = False
            sld.slide_show_transition.advance_after_time = time_in_seconds * 1000
            sld.slide_show_transition.type = slides.slideshow.TransitionType.FADE

        #set animation property for MEDIA_PLAY
        for sld in presentation.slides:
            seq = sld.timeline.main_sequence
            for e in seq:
                if (e.type == 56):   # animation.EffectType.MEDIA_PLAY
                    e.timing.trigger_type = animation.EffectTriggerType.WITH_PREVIOUS
                    e.timing.trigger_delay_time = 0

        # write the PPTX file to disk
        presentation.save(pptx_file_name + '_tmp.pptx', slides.export.SaveFormat.PPTX)

def Remove_WaterMark():
    # Load the PowerPoint presentation
    presentation = Presentation(pptx_file_name + '_tmp.pptx')
    i = 0
    # Iterate over the shapes on the slide
    for slide in presentation.slides:
        for shape in slide.shapes:
            # Check if the shape name starts with 'P' followed by a digit
            #print(shape.name)
            if shape.name == "TextBox":        
                #print("removeing... " + shape.name)
                i += 1
                shape.element.getparent().remove(shape.element)                
    # Save the modified presentation
    #print(f'Total {i} TextBox removed')
    presentation.save(pptx_file_name + '_autoplay.pptx')  

if os.path.exists('audio\\') == False:
    os.makedirs('audio\\')

print("Start adding audio to PPTX...")
Add_Audio_to_PPTX()
Remove_WaterMark()
print("Done! Check: " + pptx_file_name + '_autoplay.pptx !')